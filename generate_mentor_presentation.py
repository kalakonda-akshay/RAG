"""
Generates an exhaustive, professional PowerPoint (.pptx) presentation deck
for presenting the Orbit Offline Multimodal RAG Platform to a Mentor/Evaluator.
"""
import os
import sys

PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
REPORTS_DIR = os.path.join(PROJECT_DIR, "data", "reports")
os.makedirs(REPORTS_DIR, exist_ok=True)

PPTX_PATH = os.path.join(REPORTS_DIR, "Orbit_Offline_RAG_Mentor_Presentation.pptx")

try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    pptx = None


def create_mentor_presentation():
    if pptx is None:
        print("python-pptx not installed, skipping presentation creation.")
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG_DARK = RGBColor(11, 13, 20)       # #0B0D14
    COLOR_PRIMARY = RGBColor(99, 102, 241)     # #6366F1
    COLOR_SECONDARY = RGBColor(168, 85, 247)   # #A855F7
    COLOR_TEXT_LIGHT = RGBColor(243, 244, 246) # #F3F4F6
    COLOR_TEXT_MUTED = RGBColor(156, 163, 175)# #9CA3AF
    COLOR_CARD_BG = RGBColor(20, 23, 34)       # #141722
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_GREEN = RGBColor(52, 211, 153)      # #34D399

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="ORBIT PLATFORM &middot; SYSTEM PRESENTATION"):
        # Header Box
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0

        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_SECONDARY

        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_PRIMARY):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # SLIDE 1: Title Slide
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG_DARK
    bg1.line.fill.background()

    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "🧠 ORBIT PLATFORM"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY

    p2 = tf1.add_paragraph()
    p2.text = "Offline Multimodal Knowledge Engine & AI Assistant"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_TEXT_LIGHT

    p3 = tf1.add_paragraph()
    p3.text = "Complete Architecture, Algorithms, Technical Stack & Implementation Guide"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_before = Pt(15)

    p4 = tf1.add_paragraph()
    p4.text = "Prepared for Mentor Review & Final Evaluation"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_GREEN
    p4.space_before = Pt(25)

    # SLIDE 2: Problem Statement & Vision
    s2 = prs.slides.add_slide(blank_layout)
    bg2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = COLOR_BG_DARK
    bg2.line.fill.background()
    add_header(s2, "1. Executive Summary: Problem Statement & Vision")

    c1 = add_card(s2, 0.8, 1.6, 5.6, 5.2)
    tb_c1 = s2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "⚠️ The Cloud AI Privacy Problem"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_SECONDARY
    
    bullets_c1 = [
        ("Data Leakage Risk: ", "Sending corporate contracts, financial data, or patient records to third-party cloud APIs poses extreme compliance risks (GDPR, HIPAA)."),
        ("Recurring API Costs: ", "Pay-per-token pricing models scale unpredictably for large document processing."),
        ("Cloud Dependency: ", "Requires steady internet connection; unavailable in air-gapped or remote deployment environments."),
    ]
    for bold_prefix, text in bullets_c1:
        p = tf_c1.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run(); r1.text = bold_prefix; r1.font.bold = True; r1.font.size = Pt(12); r1.font.color.rgb = COLOR_WHITE
        r2 = p.add_run(); r2.text = text; r2.font.size = Pt(11); r2.font.color.rgb = COLOR_TEXT_MUTED

    c2 = add_card(s2, 6.8, 1.6, 5.6, 5.2)
    tb_c2 = s2.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "💡 The Orbit Solution"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
    
    bullets_c2 = [
        ("100% Air-Gapped Operation: ", "Runs entirely on local hardware with zero external network transmission."),
        ("Multimodal Format Ingestion: ", "Native support for PDF, Word, PPTX, Excel/CSV, Images, Audio, Video, HTML, and SQLite."),
        ("Sub-Second Local Retrieval: ", "Hybrid BM25 + Dense vector search delivering instant citations."),
        ("Zero Operating Cost: ", "Uses local open-source LLMs (llama3.2:3b, nomic-embed-text) with zero API charges."),
    ]
    for bold_prefix, text in bullets_c2:
        p = tf_c2.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run(); r1.text = bold_prefix; r1.font.bold = True; r1.font.size = Pt(12); r1.font.color.rgb = COLOR_WHITE
        r2 = p.add_run(); r2.text = text; r2.font.size = Pt(11); r2.font.color.rgb = COLOR_TEXT_MUTED

    # SLIDE 3: System Architecture (3-Column UI)
    s3 = prs.slides.add_slide(blank_layout)
    bg3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg3.fill.solid(); bg3.fill.fore_color.rgb = COLOR_BG_DARK; bg3.line.fill.background()
    add_header(s3, "2. System Architecture: 3-Column Orbit UI")

    cols_data = [
        ("Column 1: AppSidebar", "Navigation & Controls", [
            "User Profile Header (Template User)",
            "AI Persona Switcher (5 Roles)",
            "Local LLM Model Dropdown",
            "Agentic & Self-RAG Toggles",
            "Multimodal File Drag-and-Drop",
            "Auto-Clear Upload Checkbox",
            "Hardware CPU & RAM Monitor"
        ], 0.8),
        ("Column 2: Inbox Pane", "Q&A & Knowledge List", [
            "Live Search Bar for Conversations",
            "Auto-Labeling Status Indicator",
            "Color-Coded File Extension Badges (PDF, Excel, Media, PPTX)",
            "Indexed Document Cards",
            "Date-Stamped Q&A History Stream"
        ], 4.8),
        ("Column 3: Workspace", "Q&A & 12 Tool Tabs", [
            "Orbit Platform Gradient Header",
            "Document Focus Search Filter (🎯)",
            "Chat Workspace with Citations",
            "Voice Call Mode (STT/TTS)",
            "12 Interactive Feature Tabs",
            "Speed Badge (⚡ 0.85s response)",
            "Read Aloud Speech Synthesizer"
        ], 8.8)
    ]

    for title, subtitle, items, left in cols_data:
        add_card(s3, left, 1.6, 3.7, 5.2)
        tb = s3.shapes.add_textbox(Inches(left + 0.15), Inches(1.8), Inches(3.4), Inches(4.8))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY
        p_sub = tf.add_paragraph(); p_sub.text = subtitle; p_sub.font.size = Pt(11); p_sub.font.color.rgb = COLOR_SECONDARY
        for it in items:
            p_i = tf.add_paragraph()
            p_i.space_before = Pt(6)
            p_i.text = f"• {it}"
            p_i.font.size = Pt(10)
            p_i.font.color.rgb = COLOR_TEXT_MUTED

    # SLIDE 4: Multimodal Ingestion Pipeline
    s4 = prs.slides.add_slide(blank_layout)
    bg4 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg4.fill.solid(); bg4.fill.fore_color.rgb = COLOR_BG_DARK; bg4.line.fill.background()
    add_header(s4, "3. Multimodal Ingestion & Processing Pipeline (15+ Formats)")

    ingest_cards = [
        ("PDF & Office Docs", "PyMuPDF (fitz) + python-docx + python-pptx", "Extracts structured text, headers, slides, speaker notes, and shape text. Automatic Tesseract OCR fallback for scanned pages."),
        ("Spreadsheet Data", "pandas + openpyxl", "Converts Excel (.xlsx/.xls) and CSV sheets into Markdown tables to preserve row-column relationships for accurate table RAG."),
        ("Image Intelligence", "Pillow (PIL) + Tesseract OCR", "Non-blocking image ingestion extracting text from diagrams, screenshots, and scanned files with fallback chunking."),
        ("Audio & Video Transcripts", "faster-whisper (CTranslate2) + imageio-ffmpeg", "Extracts 16kHz mono audio from video files (.mp4/.mkv) and transcribes audio (.mp3/.wav) 4x faster than real-time."),
        ("Web Archiver & SQL DBs", "bs4 + requests + sqlite3", "Crawls external webpage URLs into offline text, and registers SQLite database schemas for natural language SQL execution.")
    ]

    top_pos = 1.6
    for title, tech, desc in ingest_cards:
        add_card(s4, 0.8, top_pos, 11.7, 0.95)
        tb = s4.shapes.add_textbox(Inches(1.0), Inches(top_pos + 0.05), Inches(11.3), Inches(0.85))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = f"{title}  —  "
        p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY
        r = p.add_run(); r.text = f"Tech: {tech}"; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = COLOR_SECONDARY
        p_desc = tf.add_paragraph(); p_desc.text = desc; p_desc.font.size = Pt(10.5); p_desc.font.color.rgb = COLOR_TEXT_MUTED
        top_pos += 1.05

    # SLIDE 5: Retrieval-Augmented Generation Architecture
    s5 = prs.slides.add_slide(blank_layout)
    bg5 = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg5.fill.solid(); bg5.fill.fore_color.rgb = COLOR_BG_DARK; bg5.line.fill.background()
    add_header(s5, "4. Advanced RAG Architecture: Retrieval to Citation Generation")

    steps_rag = [
        ("Step 1: Agentic Query Decomposition", "Splits complex multi-part user questions into focused sub-queries using LLM DAG reasoning."),
        ("Step 2: Dense & Sparse Candidate Search", "Retrieves top candidates via 768-dim Dense Vector Cosine Similarity and Sparse BM25 Keyword Matching."),
        ("Step 3: Reciprocal Rank Fusion (RRF)", "Merges and reranks vector and keyword search ranks using RRF = 1 / (60 + r_vec) + 1 / (60 + r_bm25)."),
        ("Step 4: Persona Prompt & Context Synthesis", "Combines top chunks, chat memory, and selected AI persona instructions into a grounded context window."),
        ("Step 5: Self-RAG Critic Verification", "A Critic Agent evaluates generated answer against context chunks to eliminate hallucinations before output."),
        ("Step 6: Citation & GraphRAG Rendering", "Displays response with bracketed source badges ([1] doc.pdf · p.2) and extracted entity triples.")
    ]

    top_pos = 1.6
    for title, desc in steps_rag:
        add_card(s5, 0.8, top_pos, 11.7, 0.78)
        tb = s5.shapes.add_textbox(Inches(1.0), Inches(top_pos + 0.05), Inches(11.3), Inches(0.7))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(12.5); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
        p_desc = tf.add_paragraph(); p_desc.text = desc; p_desc.font.size = Pt(10.5); p_desc.font.color.rgb = COLOR_TEXT_MUTED
        top_pos += 0.85

    # SLIDE 6: Complete Technology Stack
    s6 = prs.slides.add_slide(blank_layout)
    bg6 = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg6.fill.solid(); bg6.fill.fore_color.rgb = COLOR_BG_DARK; bg6.line.fill.background()
    add_header(s6, "5. Complete Technology Stack & Rationale")

    tech_grid = [
        ("Streamlit 1.39+", "Frontend Framework", "Fast reactive UI with custom CSS injection for 3-column layout."),
        ("Ollama (llama3.2:3b)", "Local LLM Inference Engine", "Sub-second 3B parameter local LLM execution on CPU."),
        ("Ollama (nomic-embed-text)", "Dense Vector Embedder", "768-dimensional dense vector embeddings for technical retrieval."),
        ("ChromaDB Persistent Client", "Local Vector Store", "Zero-server, SQLite-backed HNSW vector index on local disk."),
        ("Rank-BM25 (BM25Okapi)", "Sparse Keyword Search", "Exact keyword matching for technical terms, part numbers, and acronyms."),
        ("PyMuPDF & Tesseract", "PDF & OCR Processing", "High-speed C-backed PDF text extraction with OCR fallback."),
        ("faster-whisper (CTranslate2)", "Audio Speech Transcription", "Int8-quantized C++ Whisper engine running 4x faster than real-time."),
        ("PyInstaller 6.22", "Standalone Executable", "Bundles Python runtime and DLLs into a single 1-click executable.")
    ]

    for idx, (name, role, reason) in enumerate(tech_grid):
        col_idx = idx % 2
        row_idx = idx // 2
        left = 0.8 + col_idx * 5.9
        top = 1.6 + row_idx * 1.3
        add_card(s6, left, top, 5.7, 1.18)
        tb = s6.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.05), Inches(5.4), Inches(1.08))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = f"{name}  "
        p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY
        r = p.add_run(); r.text = f"({role})"; r.font.size = Pt(10); r.font.color.rgb = COLOR_SECONDARY
        p_r = tf.add_paragraph(); p_r.text = reason; p_r.font.size = Pt(10); p_r.font.color.rgb = COLOR_TEXT_MUTED

    # SLIDE 7: Algorithms & Mathematical Formulations
    s7 = prs.slides.add_slide(blank_layout)
    bg7 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg7.fill.solid(); bg7.fill.fore_color.rgb = COLOR_BG_DARK; bg7.line.fill.background()
    add_header(s7, "6. Core Algorithms & Mathematical Formulations")

    algos = [
        ("1. Reciprocal Rank Fusion (RRF)", "RRF_Score(d) = Σ [ 1 / (60 + r_m(d)) ]", "Combines sparse BM25 ranks and dense vector similarity ranks into a single score without requiring score normalization across different scoring scales."),
        ("2. Okapi BM25 Keyword Scoring", "Score(D,Q) = Σ IDF(q_i) * [ f(q_i,D)*(k1+1) ] / [ f(q_i,D) + k1*(1 - b + b*(|D|/avgdl)) ]", "Calculates term frequency saturation (k1=1.5) and document length normalization (b=0.75) for precise keyword matching."),
        ("3. Dense Vector Cosine Similarity", "Cosine_Similarity(A,B) = (A • B) / (||A|| * ||B||)", "Measures the geometric angle between 768-dimensional query and document vectors to capture semantic meaning regardless of length."),
        ("4. Overlapping Window Chunking", "Window Size = 300 words | Overlap = 50 words", "Sliding window algorithm splitting documents into chunks while preserving contextual boundaries across sentence edges.")
    ]

    top_pos = 1.6
    for title, math_eq, desc in algos:
        add_card(s7, 0.8, top_pos, 11.7, 1.18)
        tb = s7.shapes.add_textbox(Inches(1.0), Inches(top_pos + 0.05), Inches(11.3), Inches(1.08))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = f"{title}  —  "
        p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY
        r = p.add_run(); r.text = math_eq; r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = COLOR_GREEN
        p_desc = tf.add_paragraph(); p_desc.text = desc; p_desc.font.size = Pt(10.5); p_desc.font.color.rgb = COLOR_TEXT_MUTED
        top_pos += 1.3

    # SLIDE 8: AI Personas & Role Switcher
    s8 = prs.slides.add_slide(blank_layout)
    bg8 = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg8.fill.solid(); bg8.fill.fore_color.rgb = COLOR_BG_DARK; bg8.line.fill.background()
    add_header(s8, "7. AI Personas & Role Adaptation Engine")

    personas_info = [
        ("💬 General Assistant", "Default Everyday Role", "Balanced, concise, direct answers with standard citation markers for general document reading."),
        ("🛠️ Technical Auditor", "Engineering & Code Focus", "Analyzes code logic, system architecture, API schemas, security risks, and performance bottlenecks."),
        ("📊 Financial Analyst", "Numbers & Margins Focus", "Focuses on revenue trends, EBITDA margins, balance sheets, and metrics formatted in Markdown tables."),
        ("⚖️ Legal Counsel", "Contracts & Compliance Focus", "Inspects contractual obligations, indemnification clauses, liability terms, and compliance risks."),
        ("🔬 Research Scientist", "Academic & Empirical Focus", "Evaluates research methodologies, experimental findings, statistical evidence, and dataset metrics.")
    ]

    top_pos = 1.6
    for title, role_type, desc in personas_info:
        add_card(s8, 0.8, top_pos, 11.7, 0.95)
        tb = s8.shapes.add_textbox(Inches(1.0), Inches(top_pos + 0.05), Inches(11.3), Inches(0.85))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = f"{title}  "
        p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY
        r = p.add_run(); r.text = f"[{role_type}]"; r.font.size = Pt(11); r.font.color.rgb = COLOR_SECONDARY
        p_desc = tf.add_paragraph(); p_desc.text = desc; p_desc.font.size = Pt(10.5); p_desc.font.color.rgb = COLOR_TEXT_MUTED
        top_pos += 1.05

    # SLIDE 9: Interactive Features (12 Tabs)
    s9 = prs.slides.add_slide(blank_layout)
    bg9 = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg9.fill.solid(); bg9.fill.fore_color.rgb = COLOR_BG_DARK; bg9.line.fill.background()
    add_header(s9, "8. Interactive Features & 12 Workspace Tabs")

    tabs_info = [
        ("💬 Chat Workspace", "Multi-turn RAG chat stream with citations, GraphRAG triples, speed metrics, and Read Aloud TTS."),
        ("🎙️ Voice Call Mode", "Microphone audio input (st.audio_input), Whisper transcription, LLM generation, and auto-speech playback."),
        ("🌐 Web Archiver", "Crawls external website URLs into clean offline text and indexes them into ChromaDB."),
        ("🗺️ Mind Map Generator", "Renders hierarchical Mermaid.js visual mind maps of document collection structure."),
        ("🧩 Topic Clusters", "Automatically groups document chunks into semantic topic clusters with descriptions."),
        ("🌐 Offline Translator", "Translates answers and summaries into 8+ languages using local LLM prompting."),
        ("⚔️ Model Battle", "Runs two local models (Model A vs Model B) side-by-side on the same prompt for output comparison."),
        ("🎓 Quiz Generator", "Generates interactive multiple-choice quiz questions based on document context."),
        ("🗄️ SQL DB Engine", "Translates natural language questions to SQLite SQL queries and executes interactive dataframes."),
        ("🛡️ PII & PDF Tools", "Scrubs PII (emails, phone numbers, SSNs) and merges PDF files offline."),
        ("🔍 Document Inspector", "Inspects raw chunk metadata, page numbers, and vector embeddings."),
        ("📄 Document Compare", "Compares clauses and text side-by-side between any two uploaded files.")
    ]

    for idx, (t_title, t_desc) in enumerate(tabs_info):
        col_idx = idx % 3
        row_idx = idx // 3
        left = 0.8 + col_idx * 3.95
        top = 1.6 + row_idx * 1.3
        add_card(s9, left, top, 3.8, 1.18)
        tb = s9.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.05), Inches(3.6), Inches(1.08))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = t_title; p.font.size = Pt(11.5); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
        p_desc = tf.add_paragraph(); p_desc.text = t_desc; p_desc.font.size = Pt(9.5); p_desc.font.color.rgb = COLOR_TEXT_MUTED

    # SLIDE 10: Automated Testing & Verification
    s10 = prs.slides.add_slide(blank_layout)
    bg10 = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg10.fill.solid(); bg10.fill.fore_color.rgb = COLOR_BG_DARK; bg10.line.fill.background()
    add_header(s10, "9. Automated System Verification & Audit Results")

    c_audit = add_card(s10, 0.8, 1.6, 11.7, 5.2)
    tb_audit = s10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.8))
    tf_audit = tb_audit.text_frame; tf_audit.word_wrap = True

    p = tf_audit.paragraphs[0]
    p.text = "🎯 Full Test Harness Score: 15 / 15 PASSED (100% SUCCESS)"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLOR_GREEN

    p_sub = tf_audit.add_paragraph()
    p_sub.text = "Automated test harness (test_all_features.py) verified 100% operational status across all modules:"
    p_sub.font.size = Pt(12); p_sub.font.color.rgb = COLOR_TEXT_LIGHT; p_sub.space_before = Pt(8)

    test_modules = [
        "1. System Monitor (get_system_stats) — PASSED",
        "2. Personas & Prompt Templates — PASSED",
        "3. PII Redaction Engine — PASSED",
        "4. Text Parser & Chunker — PASSED",
        "5. Image Parser & OCR — PASSED",
        "6. Excel & CSV Table Parser — PASSED",
        "7. Hybrid Vector BM25 Store — PASSED",
        "8. RAG Engine (generate_answer) — PASSED",
        "9. Agentic Multi-Hop & GraphRAG — PASSED",
        "10. Self-RAG Critic Agent — PASSED",
        "11. TTS Engine (Read Aloud) — PASSED",
        "12. Mind Map Generator — PASSED",
        "13. Topic Cluster Discovery — PASSED",
        "14. Session Manager (Save/Load) — PASSED",
        "15. Report & PPTX Generator — PASSED"
    ]

    for idx, tm in enumerate(test_modules):
        col_i = idx % 2
        row_i = idx // 2
        if col_i == 0:
            p_m = tf_audit.add_paragraph()
            p_m.space_before = Pt(4)
            r = p_m.add_run(); r.text = f"{tm:<45}"; r.font.size = Pt(10.5); r.font.color.rgb = COLOR_TEXT_MUTED
        else:
            r = p_m.add_run(); r.text = f"  |  {tm}"; r.font.size = Pt(10.5); r.font.color.rgb = COLOR_TEXT_MUTED

    # SLIDE 11: Future Roadmap & Conclusion
    s11 = prs.slides.add_slide(blank_layout)
    bg11 = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg11.fill.solid(); bg11.fill.fore_color.rgb = COLOR_BG_DARK; bg11.line.fill.background()
    add_header(s11, "10. Future Roadmap & Conclusion")

    c_road = add_card(s11, 0.8, 1.6, 5.6, 5.2)
    tb_road = s11.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_road = tb_road.text_frame; tf_road.word_wrap = True
    p = tf_road.paragraphs[0]; p.text = "🚀 Next-Gen Roadmap"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY

    roadmap_items = [
        ("Native Vision LLM Integration: ", "Incorporate local LLaVA / Llama 3.2-Vision models directly into the binary bundle for visual chart understanding."),
        ("Cross-Encoder Re-ranking: ", "Add bge-reranker-large model for 2-stage retrieval re-ranking."),
        ("Sub-300ms WebRTC Voice Mode: ", "Upgrade STT/TTS to WebRTC streaming for fluid real-time voice conversations."),
        ("Native Electron Desktop App: ", "Package the frontend into an Electron desktop application wrapper.")
    ]
    for b_prefix, text in roadmap_items:
        p = tf_road.add_paragraph(); p.space_before = Pt(8)
        r1 = p.add_run(); r1.text = b_prefix; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = COLOR_WHITE
        r2 = p.add_run(); r2.text = text; r2.font.size = Pt(10); r2.font.color.rgb = COLOR_TEXT_MUTED

    c_conc = add_card(s11, 6.8, 1.6, 5.6, 5.2)
    tb_conc = s11.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_conc = tb_conc.text_frame; tf_conc.word_wrap = True
    p = tf_conc.paragraphs[0]; p.text = "🏁 Conclusion & Ready Deliverables"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_GREEN

    deliv_items = [
        ("Live Local Application: ", "Running on http://localhost:8501"),
        ("Standalone Executable Bundle: ", "OfflineRAGAssistant_Setup.zip"),
        ("Source Code Release Package: ", "OfflineRAG_v1.0.0.zip"),
        ("Full PDF Technical Manual: ", "Orbit_Offline_RAG_Platform_Documentation.pdf"),
        ("Full Word Documentation: ", "Orbit_Offline_RAG_Platform_Documentation.docx"),
        ("Automated Test Score: ", "15/15 Passed (100% Operational)")
    ]
    for b_prefix, text in deliv_items:
        p = tf_conc.add_paragraph(); p.space_before = Pt(8)
        r1 = p.add_run(); r1.text = b_prefix; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = COLOR_WHITE
        r2 = p.add_run(); r2.text = text; r2.font.size = Pt(10); r2.font.color.rgb = COLOR_TEXT_MUTED

    prs.save(PPTX_PATH)
    print(f"[SUCCESS] Mentor Presentation PPTX created: {PPTX_PATH}")


if __name__ == "__main__":
    create_mentor_presentation()
