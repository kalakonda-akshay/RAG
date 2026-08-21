"""
Handles PowerPoint PPTX slide text extraction using python-pptx.
"""
import os
import pptx


def extract_text_from_pptx(file_path: str) -> list[dict]:
    """
    Extracts text from a PowerPoint (.pptx) presentation slide by slide.
    Extracts text from all shapes, tables, and speaker notes on each slide.
    Returns a list of dicts: {"text": "<slide text>", "source": "<filename>", "page": <slide_number>, "type": "pptx"}
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    filename = os.path.basename(file_path)
    prs = pptx.Presentation(file_path)
    results = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts = []

        # Extract text from shapes and tables
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        slide_texts.append(line)

            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        slide_texts.append(" | ".join(row_texts))

        # Extract speaker notes if available
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"[Speaker Notes: {notes}]")

        full_slide_text = "\n".join(slide_texts).strip()
        if full_slide_text:
            results.append({
                "text": full_slide_text,
                "source": filename,
                "page": slide_idx,
                "type": "pptx",
            })

    return results
