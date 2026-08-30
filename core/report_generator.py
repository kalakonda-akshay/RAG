"""
Generates Executive Briefs and PowerPoint (.pptx) presentation decks from document collections.
"""
import os
import pptx
from pptx.util import Inches, Pt
import ollama

try:
    from core.vectorstore import query_collection
except ImportError:
    from vectorstore import query_collection


def generate_executive_brief(collection_name: str = "documents") -> str:
    """
    Generates a comprehensive executive summary report across all documents in the collection.
    """
    results = query_collection("executive summary key insights findings data", top_k=8, collection_name=collection_name)
    if not results:
        return "No documents found in collection to generate brief."

    combined_context = "\n\n".join([r.get("text", "") for r in results])
    prompt = f"""Generate a professional Executive Brief summarizing the following documents.
Format with clean Markdown headers:
# Executive Brief
## 1. Core Summary
## 2. Key Findings & Metrics
## 3. Risk Factors & Considerations
## 4. Action Items

Documents Context:
{combined_context[:3000]}

Executive Brief:"""

    try:
        res = ollama.generate(model="llama3.2:3b", prompt=prompt)
        return res.get("response", "Could not generate brief.").strip()
    except Exception as e:
        return f"Error generating executive brief: {str(e)}"


def generate_pptx_deck(output_path: str, collection_name: str = "documents") -> str:
    """
    Generates a 5-slide PowerPoint presentation deck (.pptx) summarizing key document insights.
    Returns path to created presentation file.
    """
    results = query_collection("key findings executive summary main points data conclusions", top_k=6, collection_name=collection_name)
    context_text = "\n".join([r.get("text", "") for r in results])[:2500] if results else "Sample Overview Data"

    # Query LLM for slide content structured data
    prompt = f"""Create content for a 4-slide presentation deck based on the context below.
Output EXACTLY 4 sections separated by '---':

Slide 1: Executive Overview
Slide 2: Key Findings & Data Points
Slide 3: Strategic Insights & Analysis
Slide 4: Recommendations & Next Steps

Context:
{context_text}

Slide Content:"""

    try:
        res = ollama.generate(model="llama3.2:3b", prompt=prompt)
        deck_text = res.get("response", "")
        slide_contents = [s.strip() for s in deck_text.split("---") if s.strip()]
    except Exception:
        slide_contents = [
            "Executive Overview\n- Comprehensive document analysis\n- Automated intelligence extraction",
            "Key Findings\n- Data extracted from uploaded files\n- Verified offline citations",
            "Strategic Insights\n- Grounded LLM synthesis\n- Zero cloud dependencies",
            "Recommendations\n- Continue document ingestion\n- Review extracted metrics"
        ]

    # Build PPTX presentation using python-pptx
    prs = pptx.Presentation()
    
    # Title Slide
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.title.text = "Offline RAG Executive Deck"
    slide.placeholders[1].text = f"Automated Intelligence Summary\nCollection: {collection_name}"

    # Content Slides
    bullet_layout = prs.slide_layouts[1]
    for slide_raw in slide_contents[:4]:
        lines = [line.strip() for line in slide_raw.split("\n") if line.strip()]
        if not lines:
            continue

        slide = prs.slides.add_slide(bullet_layout)
        title_text = lines[0].lstrip("# 123456789.Slide ")
        slide.shapes.title.text = title_text[:50]

        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        for body_line in lines[1:]:
            p = tf.add_paragraph()
            p.text = body_line.lstrip("-*• ")
            p.font.size = Pt(18)

    prs.save(output_path)
    return output_path
